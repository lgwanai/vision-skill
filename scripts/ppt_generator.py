#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT Generator — 基于 gpt-image2-ppt-skills 完整方法论的幻灯片生成引擎。

核心能力:
  1. slide_spec 精确元素控制 (type / content / position / style / color)
  2. 版本链管理 (metadata.json) — generate → edit → rollback
  3. 外部真实图片贴入 (slot规划 → skeleton → 生成 → 叠加 → 质检)
  4. 模板克隆 (PPTX → 渲染 → Vision分析 → Profile → layout匹配)
  5. 智能布局 (基于素材属性 + 文字密度 + 模板候选区)
  6. 强制交互确认工作流 (信息不足→询问→确认→冒烟→全量)

设计原则:
  - 使用 vision-skill 已有的 IMAGE_API_* 配置, 不依赖 gpt-image-2
  - 生图前必须用户确认, 绝不跳过
  - 先冒烟(封面)再全量
"""

from __future__ import annotations

import argparse
import copy
import json
import os
import re
import sys
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

# =============================================================================
# Paths & Constants
# =============================================================================

SCRIPT_DIR = Path(__file__).parent
SKILL_ROOT = SCRIPT_DIR.parent
CWD = Path.cwd()

STYLES_DIR = SKILL_ROOT / "styles"
OUTPUT_BASE_DIR = "outputs"
METADATA_FILENAME = "metadata.json"

SLIDE_WIDTH = 1536
SLIDE_HEIGHT = 864

EXTERNAL_IMAGE_TYPES = {"external_image", "external_image_placeholder", "image_overlay"}
REFERENCE_ONLY_IMAGE_TYPES = {"image_reference", "visual_reference", "reference_image", "style_reference_image"}

# =============================================================================
# Prompt Constants
# =============================================================================

LANGUAGE_FONT_RULE = """
【强制语言与字体要求】
1. 幻灯片上所有文字必须使用简体中文, 严禁出现英文单词或句子(专有名词可保留英文)
2. 中文字体使用思源黑体(Source Han Sans)或苹方(PingFang SC), 严禁草书、艺术字
3. 标题字体粗体, 正文字体常规, 字号对比清晰
"""

ASPECT_RATIO_RULE = """
【画面比例 — 强制】16:9 横版宽屏 (landscape, widescreen), 宽度明显大于高度, 绝对不要方图或竖图。
"""


# =============================================================================
# Utility Helpers
# =============================================================================

def _truthy(value: Any) -> Optional[bool]:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    if isinstance(value, str):
        t = value.strip().lower()
        if t in {"1", "true", "yes", "y", "on"}: return True
        if t in {"0", "false", "no", "n", "off"}: return False
    return None


def _parse_bbox(value: Any) -> Optional[List[float]]:
    """Parse normalized [x, y, w, h] bbox, clamping to valid range."""
    if not isinstance(value, (list, tuple)) or len(value) != 4:
        return None
    try:
        x, y, w, h = [float(v) for v in value]
    except (TypeError, ValueError):
        return None
    if w <= 0 or h <= 0:
        return None
    x, y = max(0.0, min(1.0, x)), max(0.0, min(1.0, y))
    w, h = max(0.0, min(1.0 - x, w)), max(0.0, min(1.0 - y, h))
    return [x, y, w, h] if w > 0 and h > 0 else None


def _hex_to_rgb(value: Any, default: Tuple[int, int, int] = (247, 247, 245)) -> Tuple[int, int, int]:
    if not isinstance(value, str): return default
    t = value.strip().lstrip("#")
    if len(t) == 3: t = "".join(ch * 2 for ch in t)
    if len(t) != 6: return default
    try:
        return tuple(int(t[i:i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
    except ValueError:
        return default


def _resolve_asset_path(source: Any, output_dir: str) -> Optional[str]:
    if not source or not isinstance(source, str): return None
    p = Path(source).expanduser()
    if not p.is_absolute(): p = Path(output_dir) / p
    return str(p.resolve()) if p.is_file() else None


def _norm_rect_to_pixels(rect: List[float], w: int = SLIDE_WIDTH, h: int = SLIDE_HEIGHT) -> List[int]:
    return [int(round(rect[0] * w)), int(round(rect[1] * h)),
            int(round(rect[2] * w)), int(round(rect[3] * h))]


# =============================================================================
# Style Template Loading
# =============================================================================

def list_available_styles() -> List[Dict[str, str]]:
    """列出 styles/ 下所有可用风格。"""
    styles = []
    if not STYLES_DIR.is_dir(): return styles
    for md in sorted(STYLES_DIR.glob("*.md")):
        try:
            content = md.read_text(encoding="utf-8")
            name = md.stem
            desc = ""
            for line in content.splitlines():
                if line.startswith("## 风格名称"):
                    name = line.split("## 风格名称", 1)[-1].strip() or name
            m = re.search(r'## 风格描述\s*\n+(.+?)(?:\n##|\n\n|\Z)', content, re.DOTALL)
            if m: desc = m.group(1).strip().replace('\n', ' ')[:80]
            scenario = ""
            m2 = re.search(r'## 适用场景\s*\n+(.+?)(?:\n##|\n\n|\Z)', content, re.DOTALL)
            if m2: scenario = m2.group(1).strip().replace('\n', ' ')[:80]
            styles.append({"id": md.stem, "name": name, "description": desc, "scenario": scenario, "path": str(md)})
        except Exception:
            styles.append({"id": md.stem, "name": md.stem, "description": "", "scenario": "", "path": str(md)})
    return styles


def load_style_template(style_path: str) -> str:
    """从风格 .md 提取 '## 基础提示词模板' 部分。"""
    path = Path(style_path)
    if not path.is_absolute():
        # Try exact path, then with .md extension
        candidate = STYLES_DIR / path
        if not candidate.exists():
            candidate = STYLES_DIR / f"{path}.md"
        if candidate.exists(): path = candidate
    if not path.exists():
        raise FileNotFoundError(f"Style file not found: {style_path} (tried: {STYLES_DIR / style_path}, {STYLES_DIR / f'{style_path}.md'})")
    content = path.read_text(encoding="utf-8")
    marker = "## 基础提示词模板"
    idx = content.find(marker)
    if idx == -1:
        first = content.find("## ")
        if first == -1: return content
        nxt = content.find("\n## ", first + 3)
        return content[first + 3:].strip() if nxt == -1 else content[first + 3:nxt].strip()
    start = idx + len(marker)
    nxt = content.find("\n## ", start)
    return content[start:].strip() if nxt == -1 else content[start:nxt].strip()


def style_path_from_id(style_id: str) -> Optional[str]:
    """Resolve a style ID to its .md file path."""
    if style_id.endswith(".md") and Path(style_id).exists():
        return style_id
    candidate = STYLES_DIR / f"{style_id}.md"
    if candidate.exists(): return str(candidate)
    candidate2 = STYLES_DIR / style_id
    if candidate2.exists(): return str(candidate2)
    return None


# =============================================================================
# Element Text & Spec Helpers
# =============================================================================

def _stringify(v: Any) -> str:
    if v is None: return ""
    if isinstance(v, (list, tuple)): return "；".join(str(x).strip() for x in v if str(x).strip())
    if isinstance(v, dict): return json.dumps(v, ensure_ascii=False)
    return str(v).strip()


def _element_text(elem: Dict[str, Any]) -> str:
    """Human-facing text from a slide_spec element. Supports content/heading/body/items shapes."""
    c = _stringify(elem.get("content"))
    if c: return c
    parts = []
    for k in ("heading", "title", "label"):
        v = _stringify(elem.get(k))
        if v: parts.append(v); break
    for k in ("body", "text"):
        v = _stringify(elem.get(k))
        if v: parts.append(v); break
    items = elem.get("items") or elem.get("bullets")
    if items:
        parts.append(_stringify(items))
    return "：".join(parts[:2]) if parts else ""


def apply_spec_updates(spec: Dict[str, Any], element_updates: Dict[str, Any]) -> Dict[str, Any]:
    """Apply element-level updates to a slide_spec, returning updated copy."""
    updated = copy.deepcopy(spec)
    elements = updated.setdefault("elements", {})
    derived_keys = {"computed_bbox", "placement_region", "asset_size", "asset_ratio", "tailored_to_asset", "auto_layout_reason"}
    geometry_keys = {"position", "bbox", "source", "path", "asset", "fit", "slot", "slot_strategy", "tailor_to_asset", "anchor", "padding", "bleed", "layout_intent"}
    for eid, changes in element_updates.items():
        if eid in elements:
            elem = elements[eid]
            if isinstance(elem, dict) and str(elem.get("type", "")).lower() in EXTERNAL_IMAGE_TYPES and any(k in changes for k in geometry_keys):
                for k in derived_keys: elem.pop(k, None)
            elements[eid].update(changes)
        else:
            elements[eid] = changes
    return updated


def construct_edit_prompt(spec: Dict[str, Any], element_updates: Dict[str, Any]) -> str:
    """Construct natural-language edit prompt from old spec + element changes."""
    parts = ["在参考图基础上, 只修改以下内容, 保持其他所有元素(背景、装饰、布局、字体、颜色、大小)完全不变:"]
    elements = spec.get("elements", {})
    for eid, changes in element_updates.items():
        old = elements.get(eid, {})
        pos = old.get("position", "相应位置")
        old_c = old.get("content", "")
        new_c = changes.get("content", old_c)
        if "content" in changes and old_c and new_c != old_c:
            parts.append(f"将{pos}的文字从「{old_c}」改为「{new_c}」")
        elif "color" in changes:
            parts.append(f"将{pos}的颜色从{old.get('color','')}改为{changes['color']}")
        elif "style" in changes:
            parts.append(f"将{pos}的样式改为{changes['style']}")
        else:
            for k, v in changes.items():
                parts.append(f"修改{pos}的{k}为{v}")
    return "\n".join(parts)


# =============================================================================
# Session & Metadata (Version Chain)
# =============================================================================

def _find_sessions(base_dir: Optional[str] = None) -> List[Dict[str, Any]]:
    if base_dir is None: base_dir = str(CWD / OUTPUT_BASE_DIR)
    sessions = []
    if not os.path.isdir(base_dir): return sessions
    for entry in sorted(os.listdir(base_dir), reverse=True):
        sd = os.path.join(base_dir, entry)
        mp = os.path.join(sd, METADATA_FILENAME)
        if os.path.isfile(mp):
            try:
                meta = json.loads(Path(mp).read_text(encoding="utf-8"))
                sessions.append({"timestamp": entry, "dir": sd, "title": meta.get("title", "Untitled"), "slide_count": len(meta.get("slides", {}))})
            except Exception:
                sessions.append({"timestamp": entry, "dir": sd, "title": "(损坏)", "slide_count": 0})
    return sessions


def _resolve_session(session_id: str) -> str:
    if os.path.isdir(session_id) and os.path.isfile(os.path.join(session_id, METADATA_FILENAME)):
        return os.path.abspath(session_id)
    cand = str(CWD / OUTPUT_BASE_DIR / session_id)
    if os.path.isdir(cand) and os.path.isfile(os.path.join(cand, METADATA_FILENAME)):
        return cand
    for s in _find_sessions():
        if s["timestamp"].startswith(session_id): return s["dir"]
    print(f"[X] Session not found: {session_id}")
    sys.exit(1)


def _load_metadata(session_dir: str) -> Dict[str, Any]:
    mp = os.path.join(session_dir, METADATA_FILENAME)
    if not os.path.isfile(mp):
        print(f"[X] No metadata.json in {session_dir}")
        sys.exit(1)
    return json.loads(Path(mp).read_text(encoding="utf-8"))


def _save_metadata(metadata: Dict[str, Any], session_dir: str) -> None:
    mp = os.path.join(session_dir, METADATA_FILENAME)
    fd, tmp = tempfile.mkstemp(suffix=".json", prefix=".metadata_", dir=session_dir)
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        os.replace(tmp, mp)
    except Exception:
        try: os.unlink(tmp)
        except OSError: pass
        raise


def _collect_slide_numbers(metadata: Dict[str, Any]) -> List[int]:
    slides_dict = metadata.get("slides", {})
    ordered = []
    seen = set()
    for raw in metadata.get("slide_order", []) or []:
        try: n = int(raw)
        except (TypeError, ValueError): continue
        if str(n) in slides_dict and n not in seen: ordered.append(n); seen.add(n)
    nums = sorted(int(k) for k in slides_dict if k.isdigit() and int(k) not in seen)
    return ordered + nums


def _get_latest_slide_spec(metadata: Dict[str, Any], slide_number: int) -> Dict[str, Any]:
    sd = metadata.get("slides", {}).get(str(slide_number))
    if not sd: return {}
    cv = sd.get("current_version", 1)
    for v in sd.get("versions", []):
        if v.get("version") == cv: return v.get("spec", {})
    return {}


def _get_version_info(metadata: Dict[str, Any], slide_number: int, version: int) -> Optional[Dict[str, Any]]:
    sd = metadata.get("slides", {}).get(str(slide_number))
    if not sd: return None
    for v in sd.get("versions", []):
        if v.get("version") == version: return v
    return None


def _init_slide_metadata(slide_number: int, page_type: str, spec: Dict[str, Any], prompt_file: str, image_path: str) -> Dict[str, Any]:
    return {"slide_number": slide_number, "page_type": page_type, "current_version": 1, "image_snapshot": image_path,
            "versions": [{"version": 1, "action": "generate", "spec": spec, "prompt_file": prompt_file, "image_snapshot": image_path}]}


def _add_slide_version(slide_data: Dict[str, Any], new_version: int, spec: Dict[str, Any], action: str,
                       image_snapshot: str, prompt_file: str, edit_instruction: str = "", reference_version: Optional[int] = None) -> None:
    entry: Dict[str, Any] = {"version": new_version, "action": action, "spec": spec, "prompt_file": prompt_file, "image_snapshot": image_snapshot}
    if edit_instruction: entry["edit_instruction"] = edit_instruction
    if reference_version is not None: entry["reference_version"] = reference_version
    slide_data.setdefault("versions", []).append(entry)
    slide_data["current_version"] = new_version
    slide_data["image_snapshot"] = image_snapshot


# =============================================================================
# External Image Slot Support
# =============================================================================

def _collect_external_image_slots(spec: Optional[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not spec: return []
    slots = []
    for eid, elem in (spec.get("elements") or {}).items():
        if not isinstance(elem, dict): continue
        if str(elem.get("type", "")).lower() not in EXTERNAL_IMAGE_TYPES: continue
        bbox = _parse_bbox(elem.get("computed_bbox") or elem.get("bbox") or elem.get("position"))
        if not bbox: continue
        slot_info = elem.get("slot") if isinstance(elem.get("slot"), dict) else {}
        slots.append({
            "id": str(eid), "type": str(elem.get("type", "")), "bbox": bbox,
            "source": elem.get("source") or elem.get("path") or elem.get("asset"),
            "fit": str(elem.get("fit", "contain")).lower(),
            "tailored_to_asset": bool(elem.get("tailored_to_asset", False)),
            "padding": float(elem.get("padding", slot_info.get("padding", 0.0)) or 0.0),
            "bleed": float(elem.get("bleed", slot_info.get("bleed", 0.0)) or 0.0),
            "fill": elem.get("fill", slot_info.get("fill", "#F7F7F5")),
            "outline": elem.get("outline", slot_info.get("outline", "#000000")),
            "draw_frame": bool(_truthy(elem.get("draw_frame")) or False),
            "outline_width": int(elem.get("outline_width", slot_info.get("outline_width", 0)) or 0),
            "skeleton_shape": str(elem.get("skeleton_shape", "corners")).lower(),
            "skeleton_outline": elem.get("skeleton_outline", "#000000"),
            "skeleton_outline_width": int(elem.get("skeleton_outline_width", 2) or 0),
            "skeleton_fill": elem.get("skeleton_fill", "transparent"),
            "skeleton_canvas_fill": elem.get("skeleton_canvas_fill", "transparent"),
            "sanitize_background": _truthy(elem.get("sanitize_background")) is True,
            "asset_ratio": elem.get("asset_ratio"),
            "asset_size": elem.get("asset_size"),
            "layout_intent": elem.get("layout_intent"),
            "auto_layout_reason": elem.get("auto_layout_reason"),
            "detail_level": elem.get("detail_level"),
        })
    return slots


def _slots_with_real_sources(slots: List[Dict[str, Any]], output_dir: str) -> List[Dict[str, Any]]:
    return [s for s in slots if _resolve_asset_path(s.get("source"), output_dir)]


def _is_transparent_color(value: Any) -> bool:
    return isinstance(value, str) and value.strip().lower() in {"none", "transparent", "no-fill"}


def _slot_inner_rect(slot: Dict[str, Any]) -> List[float]:
    x, y, w, h = slot["bbox"]
    pad = max(0.0, float(slot.get("padding") or 0.0))
    bleed = max(0.0, float(slot.get("bleed") or 0.0))
    return [max(0.0, x + pad - bleed), max(0.0, y + pad - bleed),
            max(0.001, min(1.0 - x - pad + bleed, w - 2 * pad + 2 * bleed)),
            max(0.001, min(1.0 - y - pad + bleed, h - 2 * pad + 2 * bleed))]


def _compute_image_rect_px(slot: Dict[str, Any], asset_path: str, w: int, h: int) -> List[int]:
    inner = _slot_inner_rect(slot)
    bx, by, bw, bh = _norm_rect_to_pixels(inner, w, h)
    if slot.get("fit") == "cover": return [bx, by, bw, bh]
    try:
        from PIL import Image
        with Image.open(asset_path) as im: ratio = im.width / im.height
    except Exception:
        ratio = bw / bh if bh else 1.0
    box_ratio = bw / bh if bh else ratio
    if ratio >= box_ratio:
        dw, dh = bw, int(bw / ratio)
        dl, dt = bx, by + int((bh - dh) / 2)
    else:
        dh, dw = bh, int(bh * ratio)
        dl, dt = bx + int((bw - dw) / 2), by
    return [dl, dt, dw, dh]


def _paste_asset(canvas: Any, asset_path: str, rect: List[int], fit: str = "contain") -> None:
    from PIL import Image
    l, t, dw, dh = rect
    if dw <= 0 or dh <= 0: return
    asset = Image.open(asset_path).convert("RGB")
    if fit == "cover":
        target_r = dw / dh
        ar = asset.width / asset.height
        if ar > target_r:
            nw = int(asset.height * target_r)
            asset = asset.crop(((asset.width - nw) // 2, 0, (asset.width + nw) // 2, asset.height))
        elif ar < target_r:
            nh = int(asset.width / target_r)
            asset = asset.crop((0, (asset.height - nh) // 2, asset.width, (asset.height + nh) // 2))
    canvas.paste(asset.resize((dw, dh)), (l, t))


def _trace_dir(output_dir: str, slide_number: int) -> Path:
    return Path(output_dir) / "external_image_trace" / f"slide-{slide_number:02d}"


def create_asset_reference_skeleton(output_dir: str, slide_number: int, spec: Optional[Dict[str, Any]]) -> Optional[str]:
    """Create transparent skeleton marking external image slots for gpt-image model reference."""
    slots = _slots_with_real_sources(_collect_external_image_slots(spec), output_dir)
    if not slots: return None
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        print("(!) Pillow missing, skip skeleton")
        return None

    w, h = SLIDE_WIDTH, SLIDE_HEIGHT
    canvas_fill = slots[0].get("skeleton_canvas_fill", "transparent")
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0)) if _is_transparent_color(canvas_fill) else Image.new("RGB", (w, h), _hex_to_rgb(canvas_fill, (247, 247, 245)))
    draw = ImageDraw.Draw(img)

    trace_slots = []
    for slot in slots:
        asset = _resolve_asset_path(slot.get("source"), output_dir)
        if not asset: continue
        final_rect = _compute_image_rect_px(slot, asset, w, h)
        x1, y1, x2, y2 = final_rect[0], final_rect[1], final_rect[0] + final_rect[2], final_rect[1] + final_rect[3]
        shape = str(slot.get("skeleton_shape", "corners")).lower()
        outline = _hex_to_rgb(slot.get("skeleton_outline"), (0, 0, 0))
        owidth = int(slot.get("skeleton_outline_width") or 2)
        if shape in {"fill", "filled"}:
            fill = _hex_to_rgb(slot.get("skeleton_fill"), (200, 200, 200))
            if not _is_transparent_color(slot.get("skeleton_fill")):
                draw.rectangle([x1, y1, x2, y2], fill=fill)
        elif shape in {"outline", "rectangle", "rect"} and owidth > 0:
            draw.rectangle([x1, y1, x2, y2], outline=outline, width=owidth)
        elif shape not in {"none", "off", "false"} and owidth > 0:
            # corner markers only (default)
            mark = max(14, min(x2 - x1, y2 - y1) // 10)
            for cx, cy in ((x1, y1), (x2, y1), (x1, y2), (x2, y2)):
                hx = 1 if cx == x1 else -1
                hy = 1 if cy == y1 else -1
                draw.line([(cx, cy), (cx + mark * hx, cy)], fill=outline, width=owidth)
                draw.line([(cx, cy), (cx, cy + mark * hy)], fill=outline, width=owidth)
        trace_slots.append({"id": slot.get("id"), "source": asset, "final_image_rect_px": final_rect})

    ref_dir = Path(output_dir) / "references"
    ref_dir.mkdir(parents=True, exist_ok=True)
    ref_path = ref_dir / f"slide-{slide_number:02d}-asset-skeleton.png"
    img.save(ref_path)

    trace = _trace_dir(output_dir, slide_number)
    trace.mkdir(parents=True, exist_ok=True)
    img.save(trace / "step2-reference-outline-blank.png")
    return str(ref_path)


# =============================================================================
# Auto Layout — Smart External Image Slot Planning
# =============================================================================

def _slide_text_profile(spec: Dict[str, Any]) -> Dict[str, Any]:
    """Analyze text load/density of a slide for layout planning."""
    tc, bc, blocks = 0, 0, 0
    for elem in (spec.get("elements") or {}).values():
        if not isinstance(elem, dict): continue
        etype = str(elem.get("type", "")).lower()
        if etype in EXTERNAL_IMAGE_TYPES: continue
        length = len(_element_text(elem))
        if not length: continue
        blocks += 1
        if etype in {"title", "headline", "heading", "subtitle"}: tc += length
        else: bc += length
    load = tc + bc
    density = "heavy" if load >= 180 else ("medium" if load >= 90 else "light")
    return {"title_chars": tc, "body_chars": bc, "text_blocks": blocks, "text_load": load, "density": density}


def _external_asset_detail_level(spec: Dict[str, Any], elem_id: str, elem: Dict[str, Any]) -> str:
    """Classify whether a real asset needs read-level size or just visual presence."""
    explicit = str(elem.get("detail_level") or "").lower()
    if explicit in {"low", "photo", "decorative"}: return "low"
    if explicit in {"critical", "dense", "dense-diagram"}: return "critical"
    if explicit in {"high", "detailed", "readable", "diagram"}: return "high"
    ac = str(elem.get("asset_class") or "").lower()
    if ac in {"dense-diagram", "dense-chart", "architecture-diagram"}: return "critical"
    if ac in {"wide-document", "chart-or-document", "screenshot-document"}: return "high"
    ctx = " ".join(str(x) for x in [elem_id, elem.get("label", ""), elem.get("description", ""), spec.get("layout", "")]).lower()
    high_terms = {"paper", "table", "diagram", "chart", "screenshot", "code", "论文", "表格", "架构", "流程", "图表"}
    return "high" if any(t in ctx for t in high_terms) else "normal"


def _analyze_external_asset_visual(source: Any, output_dir: str) -> Dict[str, Any]:
    """Heuristic visual classification of a real image before layout planning."""
    path = _resolve_asset_path(source, output_dir)
    if not path: return {}
    try:
        from PIL import Image, ImageStat
        with Image.open(path) as im:
            ow, oh = im.size
            rgb = im.convert("RGB"); rgb.thumbnail((256, 256))
            stat = ImageStat.Stat(rgb); mean = stat.mean; stddev = stat.stddev
            pixels = list(rgb.getdata()); total = max(1, len(pixels))
            white_r = sum(1 for r, g, b in pixels if r > 238 and g > 238 and b > 238) / total
            color_std = sum(stddev) / 3.0
            gray = rgb.convert("L"); w, h = gray.size
            pix = gray.load(); step = max(1, min(w, h) // 96)
            edge_hits = sum(1 for y in range(0, h - step, step) for x in range(0, w - step, step)
                          if abs(int(pix[x + step, y]) - int(pix[x, y])) + abs(int(pix[x, y + step]) - int(pix[x, y])) > 42)
            edge_density = edge_hits / max(1, (w // step) * (h // step))
    except Exception:
        return {}
    ratio = ow / oh if oh else 1.0
    tags = []; tags.append("wide" if ratio >= 2.0 else ("portrait" if ratio <= 0.8 else "standard-ratio"))
    if white_r > 0.45: tags.append("white-background")
    if edge_density > 0.09: tags.append("line-art")
    if color_std > 58 and white_r < 0.35: tags.append("photo-like")
    if "photo-like" in tags and edge_density < 0.12: ac, dl = "photo", "normal"
    elif white_r > 0.45 and edge_density > 0.13: ac, dl = "dense-diagram", "critical"
    elif white_r > 0.35 and edge_density > 0.08: ac, dl = "chart-or-document", "high"
    else: ac, dl = "general-image", "normal"
    return {"asset_class": ac, "detail_level": dl, "asset_tags": tags, "asset_ratio": round(ratio, 6), "asset_size": [ow, oh],
            "visual_metrics": {"white_ratio": round(white_r, 4), "edge_density": round(edge_density, 4), "color_std": round(color_std, 2)}}


def _auto_external_image_region(spec: Dict[str, Any], elem_id: str, elem: Dict[str, Any],
                                asset_ratio: Optional[float], index: int, total: int) -> Tuple[List[float], str]:
    """Choose optimal region for a real image before generating the slide."""
    ratio = asset_ratio or 1.6
    tp = _slide_text_profile(spec); dl = _external_asset_detail_level(spec, elem_id, elem)
    tl = int(tp["text_load"]); density = str(tp["density"]); tc = int(tp["title_chars"])
    ctx = f"text={density} title_chars={tc} load={tl} images={total} ratio={ratio:.2f} detail={dl}"

    # Right rail family
    if total >= 3:
        cw = 0.78 / min(total, 3); col = index % 3
        return [0.11 + col * cw, 0.58, cw - 0.025, 0.26], f"auto grid; {ctx}"
    if total == 2:
        if dl in {"high", "critical"}:
            return (([0.48 + index * 0.24, 0.10, 0.22, 0.78], f"auto two readable rails; {ctx}") if ratio < 0.8
                    else [0.08 + index * 0.44, 0.54, 0.40, 0.36], f"auto two bottom panels; {ctx}")
        if ratio < 0.8: return [0.08 + index * 0.28, 0.18, 0.24, 0.64], f"auto two portraits; {ctx}"
        if density == "light": return [0.10 + index * 0.42, 0.52, 0.38, 0.34], f"auto two lower gallery; {ctx}"
        return [0.56, 0.15 + index * 0.35, 0.36, 0.29], f"auto two right column; {ctx}"

    if dl == "critical":
        if ratio >= 2.0: return [0.06, 0.48, 0.88, 0.44], f"auto critical wide bottom; {ctx}"
        if ratio < 0.8: return [0.52, 0.08, 0.42, 0.84], f"auto critical right rail; {ctx}"
        return [0.34, 0.06, 0.62, 0.88], f"auto critical right panel; {ctx}"
    if dl == "high":
        if ratio >= 2.0: return [0.08, 0.52, 0.84, 0.40], f"auto high wide bottom; {ctx}"
        if ratio < 0.8: return [0.58, 0.10, 0.36, 0.78], f"auto high right rail; {ctx}"
        return [0.46, 0.10, 0.48, 0.80], f"auto high right panel; {ctx}"

    if tl >= 180:
        if ratio >= 1.25: return [0.56, 0.20, 0.36, 0.50], f"auto text-heavy landscape; {ctx}"
        if ratio < 0.8: return [0.68, 0.17, 0.24, 0.64], f"auto text-heavy portrait; {ctx}"
        return [0.64, 0.22, 0.28, 0.42], f"auto text-heavy compact; {ctx}"
    if tl >= 90:
        return ([0.66, 0.14, 0.26, 0.68], f"auto medium portrait; {ctx}") if ratio < 0.8 else ([0.56, 0.18, 0.36, 0.52], f"auto medium right; {ctx}")
    return ([0.66, 0.14, 0.26, 0.70], f"auto light portrait; {ctx}") if ratio < 0.8 else ([0.54, 0.18, 0.38, 0.54], f"auto light right; {ctx}")


def _fit_bbox_to_asset_ratio(region: List[float], asset_ratio: float, padding: float = 0.0, anchor: str = "center") -> List[float]:
    """Fit a bbox inside a normalized region while preserving asset aspect ratio."""
    slide_aspect = 16.0 / 9.0; x, y, w, h = region
    pad = max(0.0, min(float(padding or 0.0), min(w, h) / 2.0 - 0.0005))
    imw, imh = max(0.001, w - 2 * pad), max(0.001, h - 2 * pad)
    target = max(0.001, asset_ratio / slide_aspect)
    if imw / imh >= target: ow, oh = imh * target, imh
    else: ow, oh = imw, imw / target
    ow, oh = min(w, ow + 2 * pad), min(h, oh + 2 * pad)
    a = (anchor or "center").lower()
    ox = x if "left" in a else (x + w - ow if "right" in a else x + (w - ow) / 2)
    oy = y if "top" in a else (y + h - oh if "bottom" in a else y + (h - oh) / 2)
    return [max(0.0, min(1.0, ox)), max(0.0, min(1.0, oy)), max(0.001, min(1.0 - ox, ow)), max(0.001, min(1.0 - oy, oh))]


def prepare_external_image_slots(spec: Optional[Dict[str, Any]], output_dir: str) -> Optional[Dict[str, Any]]:
    """Compute final external image bboxes before prompt/skeleton/PPTX generation."""
    if not spec or not spec.get("elements"): return spec
    prepared = copy.deepcopy(spec)
    items = [(eid, elem) for eid, elem in (prepared.get("elements") or {}).items()
             if isinstance(elem, dict) and str(elem.get("type", "")).lower() in EXTERNAL_IMAGE_TYPES]
    for idx, (eid, elem) in enumerate(items):
        source = elem.get("source") or elem.get("path") or elem.get("asset")
        if source:
            vas = _analyze_external_asset_visual(source, output_dir)
            if vas:
                elem.setdefault("asset_class", vas.get("asset_class"))
                elem.setdefault("asset_ratio", vas.get("asset_ratio"))
                elem.setdefault("asset_size", vas.get("asset_size"))
                if not elem.get("detail_level"): elem["detail_level"] = vas.get("detail_level")
        ar = elem.get("asset_ratio")
        position = _parse_bbox(elem.get("position") or elem.get("bbox"))
        if not position:
            region, reason = _auto_external_image_region(prepared, eid, elem, ar, idx, len(items))
            elem["position"] = region; elem["auto_layout_reason"] = reason; elem["layout_intent"] = "auto"
        else:
            region = position

        fit = str(elem.get("fit", "contain")).lower()
        strategy = str(elem.get("slot_strategy", "")).lower()
        tailor = _truthy(elem.get("tailor_to_asset"))
        should_tailor = (ar is not None and strategy not in {"exact", "fixed"}
                         and (tailor is True or strategy in {"fit-within", "contain-within", "asset-fit"}
                              or (tailor is None and "bbox" not in elem and fit == "contain")))
        elem["placement_region"] = region
        if should_tailor:
            pad = float(elem.get("padding", 0.0) or 0.0)
            anchor = str(elem.get("anchor", "center"))
            elem["computed_bbox"] = _fit_bbox_to_asset_ratio(region, ar, pad, anchor)
            elem["tailored_to_asset"] = True; elem["slot_strategy"] = strategy or "fit-within"
        else:
            elem["computed_bbox"] = _parse_bbox(elem.get("bbox") or elem.get("position"))
            elem["tailored_to_asset"] = False
            if strategy: elem["slot_strategy"] = strategy
    return prepared


# =============================================================================
# Template Analysis (复用 vision-skill VISION_API_* 模型)
# =============================================================================

def analyze_template_with_vision(pptx_path: Optional[str], images_dir: Optional[str],
                                  rebuild: bool = False) -> Dict[str, Any]:
    """Analyze PPTX template using vision-skill's VISION_API_* model.

    Returns TemplateProfile with global_style, layouts (each with page_type,
    summary, json_schema, reference_image), and theme.

    Caches results by SHA256 hash under <cwd>/template_cache/.
    """
    import hashlib, base64

    images = []
    if images_dir:
        p = Path(images_dir)
        if p.is_dir():
            images = sorted([str(f) for f in p.iterdir() if f.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}])

    # Check cache
    h = hashlib.sha256()
    if pptx_path and Path(pptx_path).exists():
        h.update(b"PPTX:"); h.update(Path(pptx_path).name.encode())
        with open(pptx_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""): h.update(chunk)
    for img in images: h.update(b"|IMG:"); h.update(Path(img).name.encode())
    ch = h.hexdigest()[:16]
    cache_dir = CWD / "template_cache"
    cache_dir.mkdir(parents=True, exist_ok=True)
    cache_path = cache_dir / f"{ch}.json"
    if cache_path.exists() and not rebuild:
        print(f"📦 template cache hit: {cache_path}")
        return json.loads(cache_path.read_text(encoding="utf-8"))

    if not images:
        print("(!) no template images available, returning basic profile")
        return {"version": "1", "source": Path(pptx_path).name if pptx_path else "unknown", "source_hash": ch, "global_style": "", "theme": {}, "layouts": []}

    client = _get_client()
    vision_model = client.get_vision_model()
    vision_url = client.vision_base_url
    vision_key = client.vision_api_key
    if not vision_model or not vision_url:
        print("[X] VISION_MODEL not configured — template analysis requires vision recognition model")
        return {"version": "1", "source": Path(pptx_path).name if pptx_path else "unknown", "source_hash": ch, "global_style": "", "theme": {}, "layouts": []}

    print(f"🔍 analyzing {len(images)} template pages (model={vision_model})...")

    # Build multi-image vision request directly (VisionClient.vision_recognition only supports single image)
    content_parts = []
    for img in images:
        with open(img, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("ascii")
        content_parts.append({"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}})

    prompt = f"""请基于下方 {len(images)} 张PPT模板页缩略图, 输出严格JSON(不要markdown):


{{
  "global_style": "80-200字概括整体视觉调性: 背景/主色/强调色/字体/装饰/留白/构图",
  "theme": {{ "primary": "#hex", "accent": "#hex", "background": "#hex" }},
  "layouts": [
    {{
      "id": "kebab-case短ID(如cover-large-title)",
      "page_index": 0,
      "page_type": "cover|agenda|section|content|data|quote|closing|other",
      "summary": "60-120字描述这页的视觉布局: 位置/装饰/文字层级, 不描述具体文字内容",
      "reuse_friendly": true,
      "reuse_reason": "为什么可复用/不可复用(30-60字)",
      "json_schema": {{
        "type": "object",
        "properties": {{ "title": {{ "type": "string", "minLength": 2, "maxLength": 24, "description": "标题" }}, "subtitle": {{ "type": "string", "minLength": 0, "maxLength": 50, "description": "副标题" }} }},
        "required": ["title"], "additionalProperties": false
      }}
    }}
  ]
}}

要求:
- layouts长度={len(images)}, 每页一个
- json_schema按JSON Schema规范, properties用英文key, description用中文
- reuse_friendly: cover=false, 独特场景/插画=false, 纯文字/卡片网格/通用列表=true
- 严格只返回JSON对象"""

    content_parts.append({"type": "text", "text": prompt})

    try:
        import requests as req
        payload = {
            "model": vision_model,
            "messages": [{"role": "user", "content": content_parts}],
            "temperature": 0.2, "stream": False,
        }
        headers = {"Authorization": f"Bearer {vision_key}", "Content-Type": "application/json"}
        resp = req.post(f"{vision_url}/chat/completions", headers=headers, json=payload, timeout=180)
        if resp.status_code != 200:
            raise RuntimeError(f"vision call failed status={resp.status_code}: {resp.text[:300]}")
        raw = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "")
        if not raw: raise RuntimeError("vision returned empty content")
        parsed = _parse_json_loose(raw)
    except Exception as e:
        print(f"(!) vision analysis failed: {e}, using basic profile")
        parsed = {"global_style": "", "theme": {}, "layouts": []}

    if not isinstance(parsed, dict) or "layouts" not in parsed:
        parsed = {"global_style": parsed.get("global_style", "") if isinstance(parsed, dict) else "",
                  "theme": parsed.get("theme", {}) if isinstance(parsed, dict) else {},
                  "layouts": []}

    # Inject reference_image paths
    for i, layout in enumerate(parsed.get("layouts", [])):
        if i < len(images): layout["reference_image"] = str(Path(images[i]).resolve())
        layout.setdefault("page_index", i)
        if layout.get("page_type") not in {"cover", "agenda", "section", "content", "data", "quote", "closing", "other"}:
            layout["page_type"] = "content"
        layout.setdefault("id", f"layout-{i + 1:02d}")
        if layout.get("page_type") == "cover": layout["reuse_friendly"] = False
        else: layout.setdefault("reuse_friendly", True)

    profile = {"version": "1", "source": Path(pptx_path).name if pptx_path else "unknown", "source_hash": ch,
               "global_style": parsed.get("global_style", ""), "theme": parsed.get("theme", {}),
               "layouts": parsed.get("layouts", [])}

    cache_path.write_text(json.dumps(profile, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"💾 template profile cached: {cache_path}")
    return profile


def _parse_json_loose(text: str) -> Any:
    """Robust JSON extraction from LLM output."""
    s = text.strip()
    if s.startswith("```"):
        s = re.sub(r"^```[a-zA-Z0-9_-]*\s*", "", s)
        if s.endswith("```"): s = s[:-3]
        s = s.strip()
    try: return json.loads(s)
    except Exception: pass
    for marker in ("{", "["):
        start = s.find(marker); end = s.rfind("}" if marker == "{" else "]")
        if start != -1 and end != -1 and end > start:
            try: return json.loads(s[start:end + 1])
            except Exception: pass
    return text


# =============================================================================
# Overlay Quality Review (复用 vision-skill VISION_API_* 模型)
# =============================================================================

def _review_overlay_with_vision(output_dir: str, slide_number: int) -> Optional[Dict[str, Any]]:
    """Use vision model to inspect whether real image overlays cover important content."""
    trace = _trace_dir(output_dir, slide_number)
    step3 = trace / "step3-generated-background.png"
    step4 = trace / "step4-final-overlay-preview.png"
    if not step3.is_file() or not step4.is_file(): return None

    client = _get_client()
    if not client.get_vision_model():
        print("(!) VISION_MODEL not configured, skip overlay review")
        return None

    print(f"🔎 slide {slide_number}: vision overlay review...")
    import base64
    with open(step3, "rb") as f: b64_3 = base64.b64encode(f.read()).decode("ascii")
    with open(step4, "rb") as f: b64_4 = base64.b64encode(f.read()).decode("ascii")

    vision_url = client.vision_base_url
    vision_key = client.vision_api_key
    vision_model = client.get_vision_model()

    prompt = """请比较两张图: 第一张是生成的PPT背景页, 第二张是贴入真实图片后的最终预览。
判断:
1. 真实图片是否覆盖了标题/正文/数字/图标/重要装饰?
2. 是否出现白色/浅色占位块/遮罩矩形? (白底图表/论文截图本身的白色不算)
3. 真实图片是否足够大、可读(针对论文/表格/架构图)?
4. 整体版式是否合理?

只返回JSON:
{"ok": bool, "cover_issue": "string", "white_placeholder_issue": bool, "layout_issue": "string",
 "recommendation": "string", "suggested_position": [x,y,w,h] 或 null, "confidence": 0到1}"""

    try:
        import requests as req
        content_parts = [
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_3}"}},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64_4}"}},
            {"type": "text", "text": prompt},
        ]
        payload = {"model": vision_model, "messages": [{"role": "user", "content": content_parts}], "temperature": 0.1, "stream": False}
        headers = {"Authorization": f"Bearer {vision_key}", "Content-Type": "application/json"}
        resp = req.post(f"{vision_url}/chat/completions", headers=headers, json=payload, timeout=120)
        raw = resp.json().get("choices", [{}])[0].get("message", {}).get("content", "") if resp.status_code == 200 else ""
        if not raw: raise RuntimeError(f"vision call failed status={resp.status_code}")
        review = _parse_json_loose(raw)
        if isinstance(review, list) and review: review = review[0]
        if not isinstance(review, dict): return None
    except Exception as e:
        print(f"(!) overlay review failed slide {slide_number}: {e}")
        return None

    review["slide_number"] = slide_number
    trace.mkdir(parents=True, exist_ok=True)
    (trace / "overlay_review.json").write_text(json.dumps(review, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"  {'✅' if review.get('ok') else '❌'} slide {slide_number}: ok={review.get('ok')}, placeholder={review.get('white_placeholder_issue')}")
    return review


def review_external_overlays(metadata: Dict[str, Any], output_dir: str, slide_numbers: List[int]) -> Dict[int, Dict[str, Any]]:
    """Run overlay quality review on all slides with external images."""
    reviews = {}
    for n in slide_numbers:
        spec = _get_latest_slide_spec(metadata, n)
        if not _slots_with_real_sources(_collect_external_image_slots(spec), output_dir): continue
        # Ensure trace files exist
        try:
            from PIL import Image
            img_path = Path(output_dir) / "images" / f"slide-{n:02d}.png"
            trace = _trace_dir(output_dir, n); trace.mkdir(parents=True, exist_ok=True)
            if img_path.exists() and not (trace / "step3-generated-background.png").exists():
                Image.open(img_path).convert("RGB").resize((SLIDE_WIDTH, SLIDE_HEIGHT)).save(trace / "step3-generated-background.png")
            # Create step4 overlay preview
            spec2 = prepare_external_image_slots(spec, output_dir)
            slots = _slots_with_real_sources(_collect_external_image_slots(spec2), output_dir)
            if slots and img_path.exists():
                canvas = Image.open(img_path).convert("RGB").resize((SLIDE_WIDTH, SLIDE_HEIGHT))
                for slot in slots:
                    asset = _resolve_asset_path(slot.get("source"), output_dir)
                    if asset:
                        rect = _compute_image_rect_px(slot, asset, SLIDE_WIDTH, SLIDE_HEIGHT)
                        _paste_asset(canvas, asset, rect, fit=slot.get("fit", "contain"))
                canvas.save(trace / "step4-final-overlay-preview.png")
        except Exception: pass

        r = _review_overlay_with_vision(output_dir, n)
        if r: reviews[n] = r
    return reviews


def _add_external_image_overlays(prs: Any, slide: Any, slide_number: int, spec: Dict[str, Any], output_dir: str) -> int:
    """Add real image overlays as independent PPT picture objects."""
    try:
        from pptx.dml.color import RGBColor
    except ImportError:
        return 0
    slots = _slots_with_real_sources(_collect_external_image_slots(spec), output_dir)
    added = 0
    for slot in slots:
        asset = _resolve_asset_path(slot.get("source"), output_dir)
        if not asset: continue
        x, y, w_box, h_box = slot["bbox"]
        pad = max(0.0, float(slot.get("padding") or 0.0))
        bleed = max(0.0, float(slot.get("bleed") or 0.0))
        ix = max(0.0, x + pad - bleed)
        iy = max(0.0, y + pad - bleed)
        iw = max(0.001, min(1.0 - ix, w_box - 2 * pad + 2 * bleed))
        ih = max(0.001, min(1.0 - iy, h_box - 2 * pad + 2 * bleed))
        pl, pt = int(ix * prs.slide_width), int(iy * prs.slide_height)
        pw, ph = int(iw * prs.slide_width), int(ih * prs.slide_height)
        try:
            from PIL import Image
            with Image.open(asset) as im: ratio = im.width / im.height
        except Exception:
            ratio = pw / ph if ph else 1.0
        box_ratio = pw / ph if ph else ratio
        if ratio >= box_ratio:
            dw, dh = pw, int(pw / ratio)
            dl, dt = pl, pt + int((ph - dh) / 2)
        else:
            dh, dw = ph, int(ph * ratio)
            dl, dt = pl + int((pw - dw) / 2), pt
        slide.shapes.add_picture(asset, dl, dt, width=dw, height=dh)
        added += 1
    return added


# =============================================================================
# Prompt Construction (from slide_spec or freeform content)
# =============================================================================

def _adapt_style_for_external_slots(style_template: str) -> str:
    """Remove photo/image-generation directives when real images are overlaid."""
    photo_terms = ("image", "photo", "photograph", "photography", "crop", "full-bleed", "imagery", "照片", "摄影", "图片", "图像")
    lines = []
    skip = False
    skip_sections = {"封面页构图", "内容页构图", "布局系统", "图片处理", "章节页构图", "收尾页构图"}
    for line in style_template.splitlines():
        s = line.strip()
        if re.match(r"^【(.+?)】$", s):
            skip = re.match(r"^【(.+?)】$", s).group(1) in skip_sections
            if skip: continue
        if skip: continue
        if any(t in s.lower() for t in photo_terms): continue
        lines.append(line)
    return "\n".join(lines).strip() + """
【外部真实图片后贴页覆盖规则 - 最高优先级】
- 本页真实图片会在生成后由代码精确贴入; 模型不要生成任何照片、产品照、人物照、场景摄影或类似主体
- 保留所选风格本身的配色、字体气质、结构线、留白、背景质感和抽象装饰语言
- 参考图角标区域下方只允许连续背景、轻微纹理或非常轻的抽象装饰
"""


def generate_prompt_from_spec(style_template: str, spec: Dict[str, Any], page_type: str, slide_number: int, total_slides: int, output_dir: Optional[str] = None) -> str:
    """Generate precise prompt from structured slide_spec."""
    elements = spec.get("elements", {})
    layout = spec.get("layout", "")
    external_slots = _collect_external_image_slots(spec)
    if output_dir: external_slots = _slots_with_real_sources(external_slots, output_dir)
    is_cover = page_type == "cover" or slide_number == 1
    is_data = page_type == "data" or slide_number == total_slides
    label = "封面页(cover)" if is_cover else ("数据页(data)" if is_data else "内容页(content)")
    hint = "标题/副标题为视觉焦点" if is_cover else ("突出关键数字、对比或结论" if is_data else "按层级、对齐、留白结构化呈现")

    effective_style = _adapt_style_for_external_slots(style_template) if external_slots else style_template

    element_lines = []
    for eid, elem in elements.items():
        if str(elem.get("type", "")).lower() in EXTERNAL_IMAGE_TYPES | REFERENCE_ONLY_IMAGE_TYPES: continue
        text = _element_text(elem)
        if not text and not elem.get("description"): continue
        pos = elem.get("position", "")
        style_hint = elem.get("style", "")
        color = elem.get("color", "")
        desc = elem.get("description", "")
        if desc:
            element_lines.append(f"- {eid}: {desc}")
        else:
            parts = [f"- {text}"]
            if pos: parts.append(f"位置: {pos}")
            if style_hint: parts.append(f"样式: {style_hint}")
            if color: parts.append(f"颜色: {color}")
            element_lines.append(", ".join(parts))

    elements_text = "\n".join(element_lines)
    if external_slots:
        elements_text += "\n\n【外部真实图片槽位 — 强约束】\n参考图中的细角标只标出后续代码贴入真实图片的覆盖区。不要把标题、正文、数字、图标放入这些区域。定位区底下只生成连续背景。严禁生成白色/浅色矩形、占位框或占位卡片。"

    layout_line = f"\n页面布局: {layout}" if layout else ""
    return (effective_style + "\n\n---\n\n" + f"现在请生成本组中的【{label}】, {hint}{layout_line}\n\n本页各元素的精确描述(请严格按以下布局、位置、样式生成):\n\n" + elements_text + LANGUAGE_FONT_RULE + ASPECT_RATIO_RULE)


def generate_prompt(style_template: str, page_type: str, content_text: str, slide_number: int, total_slides: int, spec: Optional[Dict[str, Any]] = None, output_dir: Optional[str] = None) -> str:
    """Generate prompt for a single slide. Uses spec if available, else freeform content."""
    if spec and spec.get("elements"):
        return generate_prompt_from_spec(style_template, spec, page_type, slide_number, total_slides, output_dir)
    is_cover = page_type == "cover" or slide_number == 1
    is_data = page_type == "data" or slide_number == total_slides
    label = "封面页(cover)" if is_cover else ("数据页(data)" if is_data else "内容页(content)")
    hint = "标题/副标题为视觉焦点" if is_cover else ("突出关键数字、对比或结论" if is_data else "按层级、对齐、留白结构化呈现")
    return (style_template + "\n\n---\n\n" + f"现在请生成本组中的【{label}】, {hint}\n本页要呈现的内容如下(请按本风格美学重新设计版式):\n\n" + content_text + LANGUAGE_FONT_RULE + ASPECT_RATIO_RULE)


# =============================================================================
# Image Generation (via vision-skill VisionClient)
# =============================================================================

def _get_client():
    sys.path.insert(0, str(SCRIPT_DIR))
    from vision_client import VisionClient
    return VisionClient()


def generate_slide_image(prompt: str, output_path: str, ref_image_path: Optional[Union[str, List[str]]] = None, client: Any = None) -> Optional[str]:
    """Call vision-skill's VisionClient to generate a slide image."""
    if client is None: client = _get_client()
    out_dir = os.path.dirname(output_path)
    if out_dir: os.makedirs(out_dir, exist_ok=True)
    print(f"  🎨 generating... prompt[:100]: {prompt[:100].replace(chr(10), ' ')}...")
    try:
        result = client.generate_image(prompt=prompt, image_urls=ref_image_path)
        if "error" in result and result["error"]:
            print(f"  ❌ API error: {result['error']}")
            return None
        image_url = None
        if isinstance(result, dict):
            if "data" in result and isinstance(result["data"], list) and result["data"]:
                image_url = result["data"][0].get("url")
            elif "url" in result: image_url = result["url"]
        if not image_url:
            print(f"  ❌ Cannot extract image URL: {str(result)[:200]}")
            return None
        import requests as req
        resp = req.get(image_url, timeout=120)
        if resp.status_code == 200:
            Path(output_path).write_bytes(resp.content)
            print(f"  ✅ saved: {output_path}")
            return output_path
        print(f"  ❌ download failed (status={resp.status_code})")
        return None
    except Exception as e:
        print(f"  ❌ generation failed: {e}")
        return None


# =============================================================================
# PPTX Packaging
# =============================================================================

def generate_pptx(output_dir: str, slide_numbers: List[int], title: str = "Untitled", metadata: Optional[Dict[str, Any]] = None) -> Optional[str]:
    """Package slide images into 16:9 .pptx with optional external image overlays."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("(!) python-pptx not installed, skip PPTX")
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    if metadata is None:
        mp = os.path.join(output_dir, METADATA_FILENAME)
        if os.path.isfile(mp):
            try: metadata = _load_metadata(output_dir)
            except SystemExit: metadata = None

    img_dir = os.path.join(output_dir, "images")
    added = 0
    overlays = 0
    for n in slide_numbers:
        ip = os.path.join(img_dir, f"slide-{n:02d}.png")
        if not os.path.exists(ip):
            print(f"  ⚠️ skip slide-{n:02d}.png (not found)")
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(ip, 0, 0, width=prs.slide_width, height=prs.slide_height)
        if metadata:
            spec = _get_latest_slide_spec(metadata, n)
            if spec: overlays += _add_external_image_overlays(prs, slide, n, spec, output_dir)
        added += 1

    if added == 0: return None
    safe = re.sub(r"[^\w一-鿿\-]+", "_", title)[:60] or "deck"
    pp = os.path.join(output_dir, f"{safe}.pptx")
    prs.save(pp)
    print(f"  📑 PPTX: {pp} ({added} slides{', ' + str(overlays) + ' overlays' if overlays else ''})")
    return pp


# =============================================================================
# Full Generation Pipeline
# =============================================================================

def run_generation(plan_path: str, style_path: str, output_dir: Optional[str] = None,
                   target_slides: Optional[List[int]] = None, concurrency: int = 3,
                   template_pptx: Optional[str] = None, template_images: Optional[str] = None,
                   template_strict: bool = False, auto_review_overlays: bool = False) -> Dict[str, Any]:
    """Execute the full slide generation pipeline with version tracking.

    Supports:
      - Built-in styles (--style dark-aurora)
      - Template cloning (--template-pptx template.pptx)
      - Auto layout for external images
      - Overlay quality review via vision model
    """
    plan = json.loads(Path(plan_path).read_text(encoding="utf-8"))
    title = plan.get("title", "Untitled Presentation")
    slides = plan.get("slides", [])
    total_slides = len(slides)

    if target_slides:
        ts = set(target_slides)
        slides = [s for s in slides if s.get("slide_number") in ts]

    # Load style (template or built-in)
    template_profile = None
    if template_pptx:
        if not template_images:
            sys.path.insert(0, str(SCRIPT_DIR))
            from render_template import render_pptx_to_pngs
            print(f"🖨️  rendering template: {template_pptx}")
            template_images = str(render_pptx_to_pngs(template_pptx))
        template_profile = analyze_template_with_vision(template_pptx, template_images)
        style_template = template_profile.get("global_style", "按模板风格生成PPT页面") if template_profile else load_style_template(style_path)
        print(f"📐 template: {template_profile.get('source', 'unknown')} ({len(template_profile.get('layouts', []))} layouts)")
    else:
        style_template = load_style_template(style_path)

    if output_dir is None:
        ts_str = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = str(CWD / OUTPUT_BASE_DIR / ts_str)
    os.makedirs(os.path.join(output_dir, "images"), exist_ok=True)

    client = _get_client()
    print(f"🎨 model: {client.get_image_model()}")
    print(f"📐 style: {style_path}")
    print(f"📄 slides: {len(slides)}/{total_slides}")
    print(f"📂 output: {output_dir}")
    if template_strict: print(f"🔒 template-strict: ON")
    print()

    # Collect tasks (skip existing)
    tasks = []
    for si in slides:
        sn = si.get("slide_number", 0)
        pt = si.get("page_type", "content")
        existing = os.path.join(output_dir, "images", f"slide-{sn:02d}.png")
        if os.path.exists(existing):
            print(f"  ⏭️ slide {sn}: exists, skip")
            continue

        raw_spec = si.get("slide_spec")
        spec = raw_spec if raw_spec and raw_spec.get("elements") else None
        # Apply auto layout for external image slots
        spec = prepare_external_image_slots(spec, output_dir)

        prompt = generate_prompt(style_template, pt, si.get("content", ""), sn, total_slides, spec=spec, output_dir=output_dir)

        # Reference images: template page (strict mode) + asset skeleton
        refs = None
        if template_strict and template_profile:
            matched = template_profile.get("layouts", [])[min(sn - 1, len(template_profile.get("layouts", [])) - 1)] if template_profile.get("layouts") else None
            if matched and matched.get("reference_image"): refs = matched["reference_image"]
        asset_ref = create_asset_reference_skeleton(output_dir, sn, spec)
        if asset_ref:
            refs = [refs, asset_ref] if refs else asset_ref

        tasks.append({"slide_number": sn, "page_type": pt, "content": si.get("content", ""), "prompt": prompt, "reference_image": refs, "slide_spec": spec})

    if not tasks:
        print("✅ all slides exist, packaging PPTX")
    else:
        from concurrent.futures import ThreadPoolExecutor, as_completed
        workers = max(1, min(concurrency, len(tasks)))
        print(f"📦 {len(tasks)} tasks, {workers} concurrent\n")
        results: Dict[int, Optional[str]] = {}

        def _run(t):
            n = t["slide_number"]
            print(f"  ▶️ [slide {n}] {t['page_type']}")
            return n, generate_slide_image(t["prompt"], os.path.join(output_dir, "images", f"slide-{n:02d}.png"), ref_image_path=t.get("reference_image"), client=client)

        with ThreadPoolExecutor(max_workers=workers) as ex:
            for fut in as_completed([ex.submit(_run, t) for t in tasks]):
                n, p = fut.result()
                results[n] = p

    # Build metadata & prompts
    existing_meta_path = os.path.join(output_dir, METADATA_FILENAME)
    if os.path.isfile(existing_meta_path):
        metadata = _load_metadata(output_dir)
        metadata["generated_at"] = datetime.now().isoformat()
    else:
        metadata = {"version": 1, "title": title, "style": style_path, "model": client.get_image_model(),
                    "generated_at": datetime.now().isoformat(), "slide_order": [], "slides": {}}

    prompts_slides = []
    import shutil as _shutil
    for si in slides:
        sn = si.get("slide_number", 0)
        pt = si.get("page_type", "content")
        ip = os.path.join(output_dir, "images", f"slide-{sn:02d}.png")
        exists = os.path.exists(ip)

        if exists and str(sn) not in metadata.get("slides", {}):
            spec = si.get("slide_spec") or {}
            prompt = generate_prompt(style_template, pt, si.get("content", ""), sn, total_slides, spec=spec, output_dir=output_dir)
            pr = f"images/slide-{sn:02d}_v0001.txt"
            ir = f"images/slide-{sn:02d}_v0001.png"
            Path(os.path.join(output_dir, pr)).parent.mkdir(parents=True, exist_ok=True)
            Path(os.path.join(output_dir, pr)).write_text(prompt, encoding="utf-8")
            _shutil.copy2(ip, os.path.join(output_dir, ir))
            metadata["slides"][str(sn)] = _init_slide_metadata(sn, pt, spec, pr, ir)
            if sn not in metadata.get("slide_order", []): metadata.setdefault("slide_order", []).append(sn)

        prompts_slides.append({"slide_number": sn, "page_type": pt, "content": si.get("content", ""),
                               "prompt": generate_prompt(style_template, pt, si.get("content", ""), sn, total_slides),
                               "image_path": ip if exists else None, "slide_spec": si.get("slide_spec")})

    prompts_data = {"metadata": {"title": title, "total_slides": total_slides, "model": client.get_image_model(),
                                  "style": style_path, "generated_at": datetime.now().isoformat()}, "slides": prompts_slides}
    Path(os.path.join(output_dir, "prompts.json")).write_text(json.dumps(prompts_data, ensure_ascii=False, indent=2), encoding="utf-8")
    _save_metadata(metadata, output_dir)

    success = sorted([si.get("slide_number", 0) for si in slides if os.path.exists(os.path.join(output_dir, "images", f"slide-{si.get('slide_number', 0):02d}.png"))])
    pptx_path = generate_pptx(output_dir, success, title, metadata) if success else None

    # Run overlay quality review if requested
    overlay_reviews = {}
    if auto_review_overlays and pptx_path:
        print("\n🔍 overlay quality review...")
        overlay_reviews = review_external_overlays(metadata, output_dir, success)
        failed = {n: r for n, r in overlay_reviews.items() if not r.get("ok")}
        if failed:
            print(f"\n⚠️  {len(failed)} slides have overlay issues:")
            for n, r in sorted(failed.items()):
                print(f"    slide {n}: {r.get('cover_issue') or r.get('recommendation', 'unknown')}")

    return {"output_dir": output_dir, "pptx_path": pptx_path, "title": title, "total_slides": total_slides,
            "generated": len(success), "failed": total_slides - len(success),
            "overlay_reviews": {str(n): r for n, r in overlay_reviews.items()}}


# =============================================================================
# Edit & Rollback
# =============================================================================

def edit_slide(session_id: str, slide_number: int, element_updates: Dict[str, Any],
               edit_instruction: str = "", edit_prompt: str = "") -> Dict[str, Any]:
    """Edit a slide with element-level precision, preserving version history."""
    import shutil as _shutil
    sd = _resolve_session(session_id)
    meta = _load_metadata(sd)
    im_dir = os.path.join(sd, "images")
    sk = str(slide_number)
    if sk not in meta.get("slides", {}):
        print(f"[X] slide {slide_number} not found")
        sys.exit(1)

    slide_data = meta["slides"][sk]
    cv = slide_data["current_version"]
    cur_spec = _get_latest_slide_spec(meta, slide_number)
    updated_spec = apply_spec_updates(cur_spec, element_updates) if element_updates else cur_spec

    edit_content = edit_prompt or (construct_edit_prompt(cur_spec, element_updates) if element_updates else "")
    if not edit_content:
        print("[X] no edit prompt or element updates")
        sys.exit(1)

    ref_path = os.path.join(sd, slide_data.get("image_snapshot", ""))
    if not os.path.isfile(ref_path): ref_path = os.path.join(im_dir, f"slide-{slide_number:02d}.png")

    nv = cv + 1
    cur_img = os.path.join(im_dir, f"slide-{slide_number:02d}.png")
    versioned = os.path.join(im_dir, f"slide-{slide_number:02d}_v{cv:04d}.png")
    if os.path.exists(cur_img): _shutil.copy2(cur_img, versioned)

    pr = f"images/slide-{slide_number:02d}_v{nv:04d}.txt"
    Path(os.path.join(sd, pr)).write_text(edit_content, encoding="utf-8")

    asset_ref = create_asset_reference_skeleton(sd, slide_number, updated_spec)
    refs = [ref_path, asset_ref] if (asset_ref and ref_path) else (asset_ref or ref_path)

    print(f"Editing slide {slide_number} (v{cv} → v{nv})")
    nimg = generate_slide_image(edit_content, os.path.join(im_dir, f"slide-{slide_number:02d}.png"), ref_image_path=refs)
    if not nimg:
        print("[X] edit generation failed")
        sys.exit(1)

    _shutil.copy2(cur_img, os.path.join(im_dir, f"slide-{slide_number:02d}_v{nv:04d}.png"))
    _add_slide_version(slide_data, nv, updated_spec, "edit", f"images/slide-{slide_number:02d}_v{nv:04d}.png", pr, edit_instruction or edit_content, cv)
    _save_metadata(meta, sd)
    pptx_path = generate_pptx(sd, _collect_slide_numbers(meta), meta.get("title", "Untitled"), meta)
    return {"slide_number": slide_number, "version": nv, "image_path": nimg, "pptx_path": pptx_path}


def rollback_slide(session_id: str, slide_number: int, to_version: int) -> Dict[str, Any]:
    """Rollback a slide to a previous version's spec, regenerating the image."""
    import shutil as _shutil
    sd = _resolve_session(session_id)
    meta = _load_metadata(sd)
    im_dir = os.path.join(sd, "images")

    target = _get_version_info(meta, slide_number, to_version)
    if not target:
        print(f"[X] version {to_version} not found")
        sys.exit(1)

    target_spec = target.get("spec", {})
    target_img = os.path.join(sd, target.get("image_snapshot", ""))
    style = meta.get("style", "")
    style_tmpl = load_style_template(style) if style and Path(style).exists() else ""

    if target_spec.get("elements"):
        prompt = generate_prompt_from_spec(style_tmpl or "按以下元素描述生成幻灯片页", target_spec, meta["slides"][str(slide_number)].get("page_type", "content"), slide_number, len(meta.get("slides", {})))
    else:
        pf = os.path.join(sd, target.get("prompt_file", ""))
        prompt = Path(pf).read_text(encoding="utf-8") if pf and os.path.isfile(pf) else "请生成一张与参考图视觉风格完全一致的幻灯片页面"

    cv = meta["slides"][str(slide_number)]["current_version"]
    nv = cv + 1
    cur_img = os.path.join(im_dir, f"slide-{slide_number:02d}.png")
    if os.path.exists(cur_img): _shutil.copy2(cur_img, os.path.join(im_dir, f"slide-{slide_number:02d}_v{cv:04d}.png"))

    pr = f"images/slide-{slide_number:02d}_v{nv:04d}.txt"
    Path(os.path.join(sd, pr)).write_text(prompt, encoding="utf-8")

    ref = target_img if os.path.isfile(target_img) else None
    print(f"Rolling back slide {slide_number} (v{cv} → v{nv}, from v{to_version})")
    nimg = generate_slide_image(prompt, os.path.join(im_dir, f"slide-{slide_number:02d}.png"), ref_image_path=ref)
    if not nimg:
        print("[X] rollback failed")
        sys.exit(1)

    _shutil.copy2(cur_img, os.path.join(im_dir, f"slide-{slide_number:02d}_v{nv:04d}.png"))
    _add_slide_version(meta["slides"][str(slide_number)], nv, target_spec, "rollback", f"images/slide-{slide_number:02d}_v{nv:04d}.png", pr, f"Rollback to v{to_version}", to_version)
    _save_metadata(meta, sd)
    pptx_path = generate_pptx(sd, _collect_slide_numbers(meta), meta.get("title", "Untitled"), meta)
    return {"slide_number": slide_number, "version": nv, "image_path": nimg, "pptx_path": pptx_path}


def ingest_pptx(pptx_path: str, session_name: Optional[str] = None) -> str:
    """Ingest external PPTX → render PNGs → create session for editing."""
    sys.path.insert(0, str(SCRIPT_DIR))
    from render_template import render_pptx_to_pngs
    import glob as _g, shutil as _shutil

    renders = render_pptx_to_pngs(pptx_path)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    sid = session_name or f"{ts}_{Path(pptx_path).stem}"
    sd = str(CWD / OUTPUT_BASE_DIR / sid)
    sim = os.path.join(sd, "images")
    os.makedirs(sim, exist_ok=True)

    pngs = sorted(_g.glob(os.path.join(renders, "page-*.png")))
    slide_order = []
    slides_meta = {}
    for i, png in enumerate(pngs, start=1):
        _shutil.copy2(png, os.path.join(sim, f"slide-{i:02d}_v0001.png"))
        _shutil.copy2(png, os.path.join(sim, f"slide-{i:02d}.png"))
        slide_order.append(i)
        slides_meta[str(i)] = {"slide_number": i, "page_type": "content", "current_version": 1,
                               "image_snapshot": f"images/slide-{i:02d}.png",
                               "versions": [{"version": 1, "action": "ingest", "spec": {"layout": "(待Agent分析填充)", "elements": {}},
                                             "prompt_file": "", "image_snapshot": f"images/slide-{i:02d}_v0001.png", "source_pptx": Path(pptx_path).name}]}

    metadata = {"version": 1, "title": Path(pptx_path).stem, "source_pptx": pptx_path, "slide_order": slide_order, "ingested_at": datetime.now().isoformat(), "slides": slides_meta}
    _save_metadata(metadata, sd)
    print(f"Session: {sd} ({len(slide_order)} slides)")
    return sd


# =============================================================================
# CLI Entry Points (called from vision_cli.py)
# =============================================================================

def handle_ppt_command(args):
    """Route PPT sub-commands."""
    cmd = args.ppt_command
    if cmd == "list-styles": _cli_list_styles(args)
    elif cmd == "plan": _cli_plan(args)
    elif cmd == "generate": _cli_generate(args)
    elif cmd == "edit": _cli_edit(args)
    elif cmd == "rollback": _cli_rollback(args)
    elif cmd == "ingest": _cli_ingest(args)
    elif cmd == "sessions": _cli_sessions(args)
    else: print(json.dumps({"error": f"Unknown ppt sub-command: {cmd}"}, ensure_ascii=False))


def _cli_list_styles(args):
    styles = list_available_styles()
    if getattr(args, "format", "text") == "json":
        print(json.dumps({"styles": styles}, ensure_ascii=False, indent=2))
    else:
        print(f"\n{'ID':<45} {'Name':<25} Scenario")
        print("-" * 100)
        for s in styles:
            print(f"{s['id']:<45} {s['name']:<25} {s.get('scenario', s.get('description', ''))[:50]}")
        print(f"\n{len(styles)} styles available")


def _cli_plan(args):
    sys.path.insert(0, str(SCRIPT_DIR))
    from md_to_plan import md_to_plan
    src = Path(args.plan_input)
    if not src.exists():
        print(json.dumps({"error": f"File not found: {args.plan_input}"}, ensure_ascii=False)); return
    dst = Path(args.output) if args.output else src.with_suffix(".json")
    plan = md_to_plan(src.read_text(encoding="utf-8"))
    dst.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"status": "ok", "input": str(src), "output": str(dst), "total_slides": plan["total_slides"]}, ensure_ascii=False))


def _cli_generate(args):
    if not args.plan or not args.style:
        print(json.dumps({"error": "需要 --plan 和 --style"}, ensure_ascii=False)); return
    sp = style_path_from_id(args.style)
    if not sp and not args.template_pptx:
        print(json.dumps({"error": f"Style not found: {args.style} (and no --template-pptx)"}, ensure_ascii=False)); return
    target = [int(x.strip()) for x in args.slides.split(",")] if args.slides else None
    result = run_generation(args.plan, sp or args.style, output_dir=args.output, target_slides=target,
                            concurrency=args.concurrency, template_pptx=args.template_pptx,
                            template_images=args.template_images, template_strict=args.template_strict,
                            auto_review_overlays=args.auto_review_overlays)
    print(json.dumps(result, ensure_ascii=False, indent=2))


def _cli_edit(args):
    if not args.session or not args.edit:
        print(json.dumps({"error": "需要 --edit 和 --session"}, ensure_ascii=False)); return
    updates = json.loads(args.element_updates) if args.element_updates else {}
    result = edit_slide(args.session, args.edit, updates, edit_instruction=args.edit_instruction or "", edit_prompt=args.edit_prompt or "")
    print(json.dumps(result, ensure_ascii=False, indent=2))


def _cli_rollback(args):
    if not args.session or not args.rollback or not args.to_version:
        print(json.dumps({"error": "需要 --rollback, --session, --to-version"}, ensure_ascii=False)); return
    result = rollback_slide(args.session, args.rollback, args.to_version)
    print(json.dumps(result, ensure_ascii=False, indent=2))


def _cli_ingest(args):
    sd = ingest_pptx(args.ingest_pptx, session_name=args.session)
    print(json.dumps({"status": "ok", "session_dir": sd}, ensure_ascii=False))


def _cli_sessions(args):
    sessions = _find_sessions(args.output)
    if not sessions:
        print("No sessions found."); return
    print(f"\n{'Session':<20} {'Slides':<8} {'Title'}")
    print("-" * 60)
    for s in sessions: print(f"{s['timestamp']:<20} {s['slide_count']:<8} {s['title']}")


# =============================================================================
# Direct CLI entry
# =============================================================================

if __name__ == "__main__":
    p = argparse.ArgumentParser(description="PPT Generator")
    sub = p.add_subparsers(dest="ppt_command")

    ls = sub.add_parser("list-styles"); ls.add_argument("--format", choices=["text", "json"], default="text")
    pl = sub.add_parser("plan"); pl.add_argument("plan_input"); pl.add_argument("-o", "--output")
    gn = sub.add_parser("generate"); gn.add_argument("--plan", required=True); gn.add_argument("--style", required=True)
    gn.add_argument("--slides"); gn.add_argument("--output"); gn.add_argument("--concurrency", type=int, default=3)
    gn.add_argument("--template-pptx", help="PPTX template to clone (enables template mode)")
    gn.add_argument("--template-images", help="Pre-rendered template page PNGs directory")
    gn.add_argument("--template-strict", action="store_true", help="Use template pages as reference images (high fidelity)")
    gn.add_argument("--auto-review-overlays", action="store_true", help="Run vision-based overlay quality review after generation")
    ed = sub.add_parser("edit"); ed.add_argument("--edit", type=int, required=True); ed.add_argument("--session", required=True)
    ed.add_argument("--element-updates"); ed.add_argument("--edit-instruction"); ed.add_argument("--edit-prompt")
    rb = sub.add_parser("rollback"); rb.add_argument("--rollback", type=int, required=True); rb.add_argument("--session", required=True); rb.add_argument("--to-version", type=int, required=True)
    ig = sub.add_parser("ingest"); ig.add_argument("--ingest-pptx", required=True); ig.add_argument("--session")
    ss = sub.add_parser("sessions"); ss.add_argument("--output")

    args = p.parse_args()
    handle_ppt_command(args)
